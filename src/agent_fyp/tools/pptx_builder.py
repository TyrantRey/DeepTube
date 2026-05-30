"""Slide generation tool: Markdown summary -> .pptx via python-pptx.

`generate_learning_path` parses the well-formed Markdown produced by the
summarizer (H1 title, bullet points under「重點摘要」, a「小結」conclusion
section) into a slide deck: a title slide, one or more bullet slides, and a
conclusion slide.
"""



import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

_BULLET_RE = re.compile(r"^\s*[-*]\s+(.*)$")
_H1_RE = re.compile(r"^#\s+(.*)$")
_H2_RE = re.compile(r"^##\s+(.*)$")

_MAX_BULLETS_PER_SLIDE = 6


class ParsedSummary:
    """The pieces of a summary needed to build a deck."""

    def __init__(self, title: str, bullets: list[str], conclusion: str) -> None:
        self.title = title
        self.bullets = bullets
        self.conclusion = conclusion


def parse_summary_markdown(summary_md: str, fallback_title: str = "影片摘要") -> ParsedSummary:
    """Extract title, key-point bullets, and conclusion from the summary Markdown."""
    title = fallback_title
    bullets: list[str] = []
    conclusion_lines: list[str] = []
    section: str | None = None

    for raw in summary_md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue

        h1 = _H1_RE.match(line)
        if h1:
            title = h1.group(1).strip() or fallback_title
            continue

        h2 = _H2_RE.match(line)
        if h2:
            heading = h2.group(1).strip()
            if "重點" in heading:
                section = "points"
            elif "小結" in heading or "結論" in heading or "總結" in heading:
                section = "conclusion"
            else:
                section = "other"
            continue

        bullet = _BULLET_RE.match(line)
        if section == "points" and bullet:
            bullets.append(bullet.group(1).strip())
        elif section == "conclusion":
            text = bullet.group(1).strip() if bullet else line.strip()
            conclusion_lines.append(text)

    return ParsedSummary(title, bullets, " ".join(conclusion_lines))


def _add_bullet_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = heading
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = bullet
        para.font.size = Pt(18)


def generate_learning_path(
    summary_md: str, out_path: str | Path, title: str | None = None
) -> str:
    """Build a .pptx deck from the summary Markdown and return its path."""
    parsed = parse_summary_markdown(summary_md, fallback_title=title or "影片摘要")
    deck_title = title or parsed.title

    prs = Presentation()

    # Title slide.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = deck_title
    if title_slide.placeholders[1] is not None:
        title_slide.placeholders[1].text = "AI YouTube 影片知識萃取助理"

    # Key-point slides (chunked).
    bullets = parsed.bullets or ["（無重點內容）"]
    for start in range(0, len(bullets), _MAX_BULLETS_PER_SLIDE):
        chunk = bullets[start : start + _MAX_BULLETS_PER_SLIDE]
        heading = "重點摘要" if start == 0 else "重點摘要 (續)"
        _add_bullet_slide(prs, heading, chunk)

    # Conclusion slide.
    if parsed.conclusion:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "小結"
        slide.placeholders[1].text_frame.text = parsed.conclusion

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return str(out_path)
