"""Stage 3 (slides) unit tests: Markdown parsing and .pptx generation."""

from pptx import Presentation

from agent_fyp.tools.pptx_builder import generate_learning_path, parse_summary_markdown


def test_parse_summary_markdown(sample_summary_md):
    parsed = parse_summary_markdown(sample_summary_md)

    assert parsed.title == "Python 入門教學"
    assert len(parsed.bullets) == 3
    assert parsed.bullets[0].startswith("[00:00]")
    assert "基礎語法" in parsed.conclusion


def test_generate_learning_path_creates_pptx(sample_summary_md, tmp_path):
    out = tmp_path / "deck.pptx"
    path = generate_learning_path(sample_summary_md, out)

    assert out.exists()
    prs = Presentation(path)
    # Title slide + 1 bullet slide (3 bullets < 6/slide) + conclusion slide.
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Python 入門教學"


def test_generate_learning_path_chunks_many_bullets(tmp_path):
    bullets = "\n".join(f"- [00:{i:02d}] point {i}" for i in range(8))
    md = f"# Big\n\n## 重點摘要\n{bullets}\n\n## 小結\nDone."
    out = tmp_path / "big.pptx"

    prs = Presentation(generate_learning_path(md, out))

    # Title + 2 bullet slides (8 bullets, 6 per slide) + conclusion = 4.
    assert len(prs.slides) == 4
